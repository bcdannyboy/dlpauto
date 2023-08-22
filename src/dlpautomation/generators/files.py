from email.mime.multipart import MIMEMultipart
import json
from docx import Document
from datetime import datetime
import os
from utils.strings import genRandomFilename
import xlsxwriter
from pptx import Presentation
from email.mime.application import MIMEApplication
from email.message import EmailMessage
import yaml
from reportlab.pdfgen import canvas
from dict2xml import dict2xml
from tabulate import tabulate
import pyzipper
from zipfile import ZipFile
import jpype
import asposediagram
jpype.startJVM()
from asposediagram.api import *
import logging

# get the output path for generated files
def genOutPath():
    logging.debug("generators.files running genOutPath()")
    project_dir = os.path.abspath(os.path.dirname(__file__))
    gendir = "generated"

    gendirpath = os.path.join(project_dir, gendir)
    if not os.path.exists(gendirpath):
        os.mkdir(gendirpath)
        logging.debug("generators.files running genOutPath() - created gendirpath = %s" % gendirpath)
        return gendirpath
    else:
        logging.debug("generators.files running genOutPath() - gendirpath = %s" % gendirpath)
        return gendirpath

# generate a word document and add NPI to it
def generateWord(Title, NPI):
    logging.debug("generators.files running generateWord()")
    if isinstance(NPI, list):
        logging.debug("generators.files running generateWord() - NPI is a list")
        NPI = '\n'.join(NPI)

    document = Document()
    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    header = f"DLP Assement Document: {Title} - {gentime}"
    document.add_heading(header)
    logging.debug("generators.files running generateWord() - added heading = %s" % header)

    p = document.add_paragraph(NPI)
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.docx")

    document.save(outpath)
    logging.debug("generators.files running generateWord() - saved document to outpath = %s" % outpath)
    return outpath

# generate an excel document and add NPI to it
def generateExcel(Title, NPI):
    logging.debug("generators.files running generateExcel()")
    if isinstance(NPI, str):
        logging.debug("generators.files running generateExcel() - NPI is a string")
        NPI = [NPI]
    if isinstance(NPI, dict):
        logging.debug("generators.files running generateExcel() - NPI is a dict")
        NPI = NPI["data"]

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    Rows = [{"Title": f"DLP Assement Document: {Title} - {gentime}"}]
    logging.debug("generators.files running generateExcel() - added title row = %s" % Rows[0])

    for i in range(len(NPI)):
        Rows.append([f"data-{i}", NPI[i]])

    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.xlsx")

    workbook = xlsxwriter.Workbook(outpath)
    worksheet = workbook.add_worksheet()

    row = 0
    for rowdat in Rows:
        if isinstance(rowdat, dict):
            first_key, first_value = list(rowdat.items())[0]
            worksheet.write(row, 0, first_key)
            worksheet.write(row, 1, first_value)
        elif isinstance(rowdat, list):
            worksheet.write(row, 0, rowdat[0])
            worksheet.write(row, 1, rowdat[1])
        else:
            worksheet.write(row, 0, "row_0")
            worksheet.write(row, 1, str(rowdat))

        row += 1

    logging.debug("generators.files running generateExcel() - saved workbook to outpath = %s" % outpath)
    return outpath

# generate a powerpoint document and add NPI to it
def generatePowerPoint(Title, NPI):
    logging.debug("generators.files running generatePowerPoint()")
    if isinstance(NPI, list):
        logging.debug("generators.files running generatePowerPoint() - NPI is a list")
        NPI = '\n'.join(NPI)
    if isinstance(NPI, dict):
        NPI = '\n'.join(NPI["data"])

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")

    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = f"DLP Assement Document: {Title} - {gentime}"
    subtitle = slide.placeholders[1]
    subtitle.text = NPI
    logging.debug("generators.files running generatePowerPoint() - added title = %s" % title.text)

    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.pptx")

    prs.save(outpath)
    logging.debug("generators.files running generatePowerPoint() - saved presentation to outpath = %s" % outpath)
    return outpath

def generateVisio(Title, NPI):
    logging.debug("generators.files running generateVisio()")
    if isinstance(NPI, list):
        logging.debug("generators.files running generateVisio() - NPI is a list")
        NPI = '\n'.join(NPI)

    diagram = Diagram()
    page = diagram.getPages().getPage(0)

    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.vsdx")

    # create a new text box shape
    shape_id = page.addShape(4.25, 5.5, 2, 1, "Text")
    shape = page.getShapes().getShape(shape_id)

    # set the text of the shape
    text = shape.getText()
    text.getValue().add("DLP Assement Document: {Title} - {gentime}\n{NPI}")

    # set the line color of the shape to none
    shape.getLine().getLineColor().setValue(ColorValue.NONE)

    # set the fill color of the shape to none
    shape.getFill().getFillForegnd().setValue(ColorValue.NONE)

    logging.debug("generators.files running generateVisio() - added text box")

    diagram.save(outpath, SaveFileFormat.VSDX)
    logging.debug("generators.files running generateVisio() - saved visio file")


# generate eml file and add NPI to it in the location or as an attachment
def generateEML(Title, Location, NPI, Attachment=None):
    logging.debug("generators.files running generateEML()")
    if isinstance(NPI, dict):
        NPI = ' - '.join(NPI["data"]).replace('\n', '<>')
    if isinstance(NPI, list):
        if Location == 'subject':
            logging.debug("generators.files running generateEML() - NPI is a list and location is subject")
            NPI = ' '.join(NPI)
        else:
            logging.debug("generators.files running generateEML() - NPI is a list and location is not subject")
            NPI = '\n'.join(NPI)

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")

    msg = MIMEMultipart()
    msg['From'] = 'dlp@dlp.com'
    msg['To'] = 'dlp@dlp.com'
    if Location == 'subject':
        msg['Subject'] = NPI
        msg.attach(EmailMessage(f"DLP Assement Document: {Title} - {gentime}"))
    elif Location == 'body':
        msg['Subject'] = f"DLP Assement Document: {Title} - {gentime}"
        msg.attach(EmailMessage(NPI))
    else:
        msg['Subject'] = f"DLP Assement Document: {Title} - {gentime}"
        msg.attach(EmailMessage(f"DLP Assement Document: {Title} - {gentime}"))

    if Attachment != None:
        with open(Attachment, 'rb') as f:
            file_name = os.path.basename(Attachment)
            file_data = f.read()
            file_mime = MIMEApplication(file_data, name=file_name)
            msg.attach(file_mime)

    logging.debug(f"generators.files running generateEML() - added title = DLP Assement Document: {Title} - {gentime}")

    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.eml")

    out = open(outpath, "w+")
    out.write(msg.as_string())

    logging.debug("generators.files running generateEML() - saved message to outpath = %s" % outpath)
    return outpath

# generate a txt or csv file and add NPI to it
def generateText(Title, NPI, csv=False, bad=False):
    logging.debug("generators.files running generateText()")
    if isinstance(NPI, dict):
        NPI = ','.join(NPI["data"]).replace("\n","<>")
    if isinstance(NPI, list):
        if csv:
            logging.debug("generators.files running generateText() - NPI is a list and csv is true")
            NPI = ','.join(NPI)
        else:
            logging.debug("generators.files running generateText() - NPI is a list and csv is false")
            NPI = '\n'.join(NPI)

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    if csv:
        NPI = f"DLP Assement Document: {Title} - {gentime}," + NPI
    else:
        NPI = f"DLP Assement Document: {Title} - {gentime}\n" + NPI

    logging.debug("generators.files running generateText() - added title = DLP Assement Document: {Title} - {gentime}")
    gendirpath = genOutPath()
    filepath = "bad"
    if not bad:
        if csv:
            logging.debug("generators.files running generateText() - csv is true and not bad")
            filepath = "csv"
        else:
            logging.debug("generators.files running generateText() - csv is false and not bad")
            filepath = "txt"
    else:
        logging.debug("generators.files running generateText() - bad is true")

    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.{filepath}")

    f = open(outpath, "w+").write(NPI)
    logging.debug("generators.files running generateText() - saved text to outpath = %s" % outpath)
    return outpath

# generate json or yaml file and add NPI
def generateJSONorYAML(Title, NPI, jsonorYaml=True):
    logging.debug("generators.files running generateJSONorYAML()")
    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.{'json' if jsonorYaml else 'yml'}")

    data = {
        "Title": f"DLP Assement Document: {Title} - {gentime}",
        "Data": NPI
    }

    logging.debug("generators.files running generateJSONorYAML() - json = %s" % jsonorYaml)
    text = ""
    if jsonorYaml:
        text = json.dumps(data)
    else:
        text = yaml.dump(data)

    f = open(outpath, "w+").write(text)
    logging.debug("generators.files running generateJSONorYAML() - saved text to outpath = %s" % outpath)
    return outpath

# generate a PDF file and add NPI
def generatePDF(Title, NPI):
    logging.debug("generators.files running generatePDF()")
    if isinstance(NPI, list):
        logging.debug("generators.files running generatePDF() - NPI is a list")
        NPI = '\n'.join(NPI)

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.pdf")

    pdf = canvas.Canvas(outpath)
    pdf.drawString(100, 750, f"DLP Assement Document: {Title} - {gentime}\n{NPI}")
    pdf.save()

    logging.debug("generators.files running generatePDF() - saved text to outpath = %s" % outpath)
    return outpath

# generates an XML file
def generateXML(Title, NPI):
    logging.debug("generators.files running generateXML()")
    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.xml")

    data = {
        "Title": f"DLP Assement Document: {Title} - {gentime}",
        "Data": NPI
    }

    text = dict2xml(data)

    f = open(outpath, "w+").write(text)
    logging.debug("generators.files running generateXML() - saved text to outpath = %s" % outpath)


    return outpath

# generate an HTML file and put NPI in it
def generateHTML(Title, NPI):
    logging.debug("generators.files running generateHTML()")
    if isinstance(NPI, str):
        logging.debug("generators.files running generateHTML() - NPI is a string")
        NPI = [NPI]

    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.html")

    html = f"""
    <html>
        <h1>DLP Assement Document: {Title} - {gentime}</h1>
        {tabulate(NPI, tablefmt='html')}
    </html>"""

    f = open(outpath, "w+").write(html)
    logging.debug("generators.files running generateHTML() - saved text to outpath = %s" % outpath)
    return outpath

# generates a ZIP file with a text file in it that contains npi. embed zips a file then zips that in another zip. if password is not empty, encrypt the zip
def generateZIP(Title, NPI, password='', embed=False):
    logging.debug("generators.files running generateZIP()")
    gentime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    gendirpath = genOutPath()
    outpath = os.path.join(gendirpath, f"{genRandomFilename()}.zip")

    file_to_zip = generateText(f"EMBEDDING: {Title} - {gentime}", NPI)

    if embed:
        logging.debug("generators.files running generateZIP() - embed is true")
        embed_path =  f"embed_{genRandomFilename()}.zip"
        embedded_zip = ZipFile(os.path.join(gendirpath, embed_path), mode="w")
        logging.debug("generators.files running generateZIP() - embedded_zip = %s" % embed_path)
        embedded_zip.write(file_to_zip)

        file_to_zip = os.path.join(gendirpath, embed_path)

    if password == '':
        logging.debug("generators.files running generateZIP() - password is empty")
        out = ZipFile(outpath, mode="w")
        out.write(file_to_zip)
    else:
        logging.debug("generators.files running generateZIP() - password is not empty")
        with pyzipper.AESZipFile(outpath, 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as out:
            out.setpassword(password.encode())
            out.write(file_to_zip)

    logging.debug("generators.files running generateZIP() - saved text to outpath = %s" % outpath)
    return outpath